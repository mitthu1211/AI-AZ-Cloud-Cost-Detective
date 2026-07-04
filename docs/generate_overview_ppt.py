from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.util import Inches, Pt


OUT = Path(__file__).with_name("AI_Cloud_Cost_Detective_Project_Overview.pptx")

NAVY = RGBColor(15, 23, 42)
SLATE = RGBColor(30, 41, 59)
BLUE = RGBColor(37, 99, 235)
CYAN = RGBColor(56, 189, 248)
GREEN = RGBColor(34, 197, 94)
AMBER = RGBColor(245, 158, 11)
RED = RGBColor(239, 68, 68)
WHITE = RGBColor(248, 250, 252)
MUTED = RGBColor(203, 213, 225)
CARD = RGBColor(241, 245, 249)
TEXT = RGBColor(15, 23, 42)


def set_bg(slide, color=NAVY):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, text, x, y, w, h, size=20, color=WHITE, bold=False, align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    if align:
        p.alignment = align
    return box


def add_title(slide, title, subtitle=None):
    add_textbox(slide, title, 0.55, 0.35, 12.2, 0.5, size=28, bold=True)
    if subtitle:
        add_textbox(slide, subtitle, 0.58, 0.88, 11.8, 0.38, size=11, color=MUTED)


def add_footer(slide, idx):
    add_textbox(slide, "AI Cloud Cost Detective", 0.55, 7.18, 4.0, 0.22, size=8, color=MUTED)
    add_textbox(slide, str(idx), 12.55, 7.18, 0.35, 0.22, size=8, color=MUTED, align=PP_ALIGN.RIGHT)


def add_card(slide, x, y, w, h, title, bullets, accent=BLUE):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD
    shape.line.color.rgb = RGBColor(226, 232, 240)
    shape.line.width = Pt(1)

    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(0.08), Inches(h))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()

    add_textbox(slide, title, x + 0.22, y + 0.16, w - 0.38, 0.32, size=13, color=TEXT, bold=True)
    top = y + 0.6
    for bullet in bullets:
        add_textbox(slide, "- " + bullet, x + 0.26, top, w - 0.45, 0.28, size=9.3, color=TEXT)
        top += 0.36
    return shape


def add_table_like(slide, x, y, rows, col_widths, row_h=0.42, header=True):
    for r, row in enumerate(rows):
        y0 = y + r * row_h
        for c, text in enumerate(row):
            x0 = x + sum(col_widths[:c])
            w = col_widths[c]
            shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x0), Inches(y0), Inches(w), Inches(row_h))
            shape.fill.solid()
            shape.fill.fore_color.rgb = SLATE if (header and r == 0) else RGBColor(248, 250, 252)
            shape.line.color.rgb = RGBColor(203, 213, 225)
            box = add_textbox(slide, text, x0 + 0.06, y0 + 0.07, w - 0.12, row_h - 0.08, size=8.4, color=WHITE if (header and r == 0) else TEXT, bold=(header and r == 0))
            box.text_frame.margin_left = 0
            box.text_frame.margin_right = 0


def add_flow_node(slide, text, x, y, w, accent=BLUE):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.62))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(239, 246, 255)
    shape.line.color.rgb = accent
    add_textbox(slide, text, x + 0.08, y + 0.16, w - 0.16, 0.28, size=10, color=TEXT, bold=True, align=PP_ALIGN.CENTER)
    return shape


def add_arrow(slide, x1, y1, x2, y2):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    conn.line.color.rgb = CYAN
    conn.line.width = Pt(2)
    conn.line.end_arrowhead = True


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

slides = []

# 1. Cover
s = prs.slides.add_slide(blank)
set_bg(s)
add_textbox(s, "AI Cloud Cost Detective", 0.75, 1.35, 11.8, 0.8, size=40, bold=True)
add_textbox(s, "Project overview, architecture, access requirements, and rollout plan", 0.8, 2.2, 10.5, 0.45, size=17, color=MUTED)
add_textbox(s, "Prepared for manager and customer review", 0.8, 2.72, 8.0, 0.35, size=12, color=CYAN, bold=True)
add_card(s, 0.85, 4.0, 3.6, 1.25, "Business Goal", ["Automate Azure cost review", "Surface actionable optimization fixes"], GREEN)
add_card(s, 4.85, 4.0, 3.6, 1.25, "Technical Goal", ["React + FastAPI application", "JWT auth, WebSocket progress, PostgreSQL history"], BLUE)
add_card(s, 8.85, 4.0, 3.6, 1.25, "Customer Outcome", ["Faster cost insights", "Clear evidence and fix commands"], AMBER)
add_footer(s, 1)

# 2. Executive overview
s = prs.slides.add_slide(blank)
set_bg(s)
add_title(s, "Executive Overview", "What the solution does and why it matters")
add_card(s, 0.75, 1.55, 3.7, 4.75, "Problem", [
    "Azure environments grow quickly and cost drift is hard to detect manually.",
    "Unused, oversized, or misconfigured resources create avoidable spend.",
    "Manual audits are slow and inconsistent across teams."
], RED)
add_card(s, 4.85, 1.55, 3.7, 4.75, "Solution", [
    "User selects an Azure Resource Group from the web UI.",
    "Backend scans Azure resources through Azure CLI.",
    "OpenAI analyzes resources and returns savings recommendations."
], BLUE)
add_card(s, 8.95, 1.55, 3.7, 4.75, "Value", [
    "Reduces analysis time from manual review to guided workflow.",
    "Produces clear issue severity, explanation, and fix command.",
    "Stores analysis history for audit and follow-up."
], GREEN)
add_footer(s, 2)

# 3. Current scope
s = prs.slides.add_slide(blank)
set_bg(s)
add_title(s, "Current Implemented Scope", "Built through prompt 5: end-to-end frontend/backend integration")
add_card(s, 0.7, 1.45, 3.9, 4.9, "Frontend", [
    "React, Vite, TypeScript, Tailwind UI.",
    "Login, signup, dashboard, report, and history pages.",
    "JWT stored in browser localStorage.",
    "Live progress tracker through WebSocket."
], CYAN)
add_card(s, 4.75, 1.45, 3.9, 4.9, "Backend", [
    "FastAPI service with protected API routes.",
    "Custom JWT auth using bcrypt and PyJWT.",
    "Azure resource group and resource scanning through Azure CLI.",
    "AI analysis orchestration and report persistence."
], BLUE)
add_card(s, 8.8, 1.45, 3.9, 4.9, "Data & Reporting", [
    "Azure PostgreSQL target schema for users and analyses.",
    "Local in-memory fallback for development testing.",
    "Analysis history and report drill-down endpoint.",
    "Issue severity, savings estimate, and fix command display."
], GREEN)
add_footer(s, 3)

# 4. Architecture
s = prs.slides.add_slide(blank)
set_bg(s)
add_title(s, "Solution Architecture", "High-level component view")
nodes = [
    ("React Frontend", 0.7, 1.6, 2.1, BLUE),
    ("FastAPI Backend", 3.25, 1.6, 2.1, CYAN),
    ("Azure CLI", 5.9, 0.95, 1.8, AMBER),
    ("OpenAI API", 5.9, 2.28, 1.8, GREEN),
    ("Azure PostgreSQL", 8.25, 1.6, 2.2, BLUE),
    ("Azure Resource Group", 10.85, 1.6, 2.0, AMBER),
]
for text, x, y, w, accent in nodes:
    add_flow_node(s, text, x, y, w, accent)
add_arrow(s, 2.8, 1.91, 3.25, 1.91)
add_arrow(s, 5.35, 1.78, 5.9, 1.25)
add_arrow(s, 7.7, 1.25, 10.85, 1.78)
add_arrow(s, 5.35, 2.04, 5.9, 2.59)
add_arrow(s, 5.35, 1.91, 8.25, 1.91)
add_textbox(s, "Protected REST APIs", 2.9, 2.18, 1.7, 0.3, size=8, color=MUTED, align=PP_ALIGN.CENTER)
add_textbox(s, "WebSocket progress stream", 0.95, 2.55, 4.0, 0.3, size=9, color=CYAN)
add_card(s, 1.0, 4.15, 3.3, 1.55, "Security Boundary", ["JWT protects scan, history, and resource group APIs.", "Secrets stay server-side."], BLUE)
add_card(s, 5.0, 4.15, 3.3, 1.55, "Integration Boundary", ["Azure CLI uses logged-in identity or service principal.", "OpenAI receives resource metadata for analysis."], AMBER)
add_card(s, 9.0, 4.15, 3.3, 1.55, "Persistence Boundary", ["PostgreSQL stores users and analysis history.", "Reports are retrieved per authenticated user."], GREEN)
add_footer(s, 4)

# 5. Request flow
s = prs.slides.add_slide(blank)
set_bg(s)
add_title(s, "End-to-End Request Flow", "User journey from login to final report")
steps = [
    "1. User signs up or logs in through React UI.",
    "2. FastAPI validates credentials and returns JWT.",
    "3. UI fetches Azure Resource Groups with Authorization header.",
    "4. User selects a Resource Group and starts analysis.",
    "5. Backend scans resources using Azure CLI.",
    "6. Backend sends resource metadata to OpenAI for analysis.",
    "7. WebSocket streams live progress to React.",
    "8. Results are stored and displayed as report and history."
]
for i, step in enumerate(steps):
    add_card(s, 0.8 + (i % 2) * 6.05, 1.25 + (i // 2) * 1.25, 5.45, 0.8, step, [], BLUE if i % 2 == 0 else CYAN)
add_footer(s, 5)

# 6. Key features
s = prs.slides.add_slide(blank)
set_bg(s)
add_title(s, "Key Features", "What the customer can use in the application")
add_card(s, 0.7, 1.35, 3.8, 4.9, "Cost Detection", [
    "Over-provisioned resources.",
    "Unused or idle resources.",
    "Misconfigured pricing tiers or settings.",
    "Storage and logging cost concerns."
], RED)
add_card(s, 4.75, 1.35, 3.8, 4.9, "User Experience", [
    "Secure login and signup.",
    "Resource group dropdown.",
    "Run Analysis workflow.",
    "Animated progress steps."
], BLUE)
add_card(s, 8.8, 1.35, 3.8, 4.9, "Report Output", [
    "Total resources scanned.",
    "Issues found and severity.",
    "Estimated savings.",
    "Fix commands in copyable code blocks."
], GREEN)
add_footer(s, 6)

# 7. Tech stack
s = prs.slides.add_slide(blank)
set_bg(s)
add_title(s, "Technology Stack", "Main implementation choices")
rows = [
    ["Layer", "Technology", "Purpose"],
    ["Frontend", "React, Vite, TypeScript, Tailwind", "User interface, routing, dashboard, reports"],
    ["Backend", "Python, FastAPI", "REST APIs, auth, analysis orchestration"],
    ["Auth", "JWT, bcrypt, PyJWT", "Custom user login and protected routes"],
    ["Azure Data", "Azure CLI", "Resource group and resource inventory scanning"],
    ["AI", "OpenAI API", "Cost analysis and optimization recommendations"],
    ["Database", "Azure PostgreSQL", "Users, analysis history, stored reports"],
    ["Realtime", "FastAPI WebSocket", "Progress messages during analysis"]
]
add_table_like(s, 0.7, 1.35, rows, [2.1, 3.6, 6.1], row_h=0.52)
add_footer(s, 7)

# 8. Access requirements
s = prs.slides.add_slide(blank)
set_bg(s)
add_title(s, "Accesses Required", "Minimum access checklist for development, demo, and production")
rows = [
    ["Area", "Required Access", "Why It Is Needed"],
    ["Azure Subscription", "Reader on target subscription or Resource Group", "List resource groups and resources"],
    ["Azure Cost", "Cost Management Reader, recommended", "Future cost APIs and stronger savings evidence"],
    ["Azure Identity", "Azure CLI login or Service Principal", "Backend needs an authenticated Azure identity"],
    ["Azure PostgreSQL", "DB connection string and schema create rights", "Create users and analyses tables"],
    ["OpenAI", "OPENAI_API_KEY with model access", "Generate cost analysis recommendations"],
    ["App Runtime", "Node.js 18+, Python 3.10+", "Run frontend and backend services"],
    ["Network", "Outbound to Azure, OpenAI, PostgreSQL", "API calls and database connectivity"],
    ["Security", "JWT_SECRET value", "Sign and validate app tokens"]
]
add_table_like(s, 0.55, 1.18, rows, [2.25, 4.1, 5.85], row_h=0.48)
add_footer(s, 8)

# 9. Config and secrets
s = prs.slides.add_slide(blank)
set_bg(s)
add_title(s, "Configuration and Secrets", "Environment variables and operational setup")
add_card(s, 0.75, 1.35, 3.7, 4.75, "Backend .env", [
    "OPENAI_API_KEY=<key>",
    "DATABASE_URL=<postgres connection string>",
    "JWT_SECRET=<strong random secret>",
    "Keep all secrets outside source control."
], BLUE)
add_card(s, 4.85, 1.35, 3.7, 4.75, "Azure Setup", [
    "Install Azure CLI on backend host.",
    "Run az login for local demo.",
    "For production, prefer service principal or managed identity.",
    "Grant least-privilege Reader access."
], AMBER)
add_card(s, 8.95, 1.35, 3.7, 4.75, "Database Setup", [
    "Provision Azure Database for PostgreSQL.",
    "Allow backend host network access.",
    "Use SSL as required by Azure PostgreSQL.",
    "Backend creates required tables at startup."
], GREEN)
add_footer(s, 9)

# 10. Demo plan
s = prs.slides.add_slide(blank)
set_bg(s)
add_title(s, "Demo Plan", "Recommended customer walkthrough")
add_card(s, 0.8, 1.2, 5.6, 5.4, "Live Demo Steps", [
    "Start backend: python -m uvicorn main:app --reload --host 127.0.0.1 --port 8000",
    "Start frontend: npm.cmd run dev",
    "Signup or login.",
    "Select Azure Resource Group.",
    "Run Analysis and watch progress.",
    "Review report: summary, issue severity, savings, and commands.",
    "Open History and revisit prior report."
], BLUE)
add_card(s, 7.0, 1.2, 5.4, 5.4, "Pre-Demo Checklist", [
    "Azure CLI installed and authenticated.",
    "Resource Group contains sample resources.",
    "OpenAI key configured and valid.",
    "PostgreSQL configured or local fallback accepted.",
    "Ports 5173 and 8000 available.",
    "Customer data handling expectations confirmed."
], GREEN)
add_footer(s, 10)

# 11. Security and governance
s = prs.slides.add_slide(blank)
set_bg(s)
add_title(s, "Security and Governance", "Controls and customer-facing considerations")
add_card(s, 0.7, 1.35, 3.8, 4.9, "Implemented Controls", [
    "JWT auth protects analysis, history, and resource group APIs.",
    "Passwords are hashed with bcrypt.",
    "History is filtered by authenticated user.",
    "Frontend sends token in Authorization header."
], GREEN)
add_card(s, 4.75, 1.35, 3.8, 4.9, "Recommended Controls", [
    "Use HTTPS in hosted environments.",
    "Move secrets to Azure Key Vault.",
    "Use managed identity or service principal.",
    "Add audit logs and token expiry/refresh policy."
], BLUE)
add_card(s, 8.8, 1.35, 3.8, 4.9, "Data Considerations", [
    "Resource metadata is sent to OpenAI.",
    "No fix command is executed automatically.",
    "Reports should be reviewed before remediation.",
    "Customer approval required for production data use."
], AMBER)
add_footer(s, 11)

# 12. Roadmap
s = prs.slides.add_slide(blank)
set_bg(s)
add_title(s, "Roadmap and Next Steps", "Path from prototype to production-ready solution")
add_card(s, 0.7, 1.35, 3.8, 4.9, "Near Term", [
    "Finalize .env and cloud database setup.",
    "Add detailed error handling and retry UX.",
    "Validate reports against multiple Azure resource types.",
    "Capture demo screenshots and sample report data."
], CYAN)
add_card(s, 4.75, 1.35, 3.8, 4.9, "Production Hardening", [
    "Use Azure managed identity.",
    "Deploy backend and frontend to Azure.",
    "Store secrets in Azure Key Vault.",
    "Add observability, audit logs, and CI/CD."
], BLUE)
add_card(s, 8.8, 1.35, 3.8, 4.9, "Future Enhancements", [
    "Integrate Azure Cost Management APIs.",
    "Add subscription-level scans.",
    "Add approval workflow for fix commands.",
    "Export reports to PDF or PowerPoint."
], GREEN)
add_footer(s, 12)

# 13. Appendix APIs
s = prs.slides.add_slide(blank)
set_bg(s)
add_title(s, "Appendix: API Surface", "Endpoints implemented in the current project")
rows = [
    ["Endpoint", "Auth", "Purpose"],
    ["POST /api/auth/signup", "No", "Create user and return JWT"],
    ["POST /api/auth/login", "No", "Validate credentials and return JWT"],
    ["GET /api/resource-groups", "Yes", "List Azure Resource Groups"],
    ["POST /api/analyze", "Yes", "Run AI cost analysis for selected Resource Group"],
    ["GET /api/history", "Yes", "List prior analyses for authenticated user"],
    ["GET /api/history/{analysis_id}", "Yes", "Open full stored report details"],
    ["WS /ws/progress/{analysis_id}", "No token currently", "Stream live progress messages"]
]
add_table_like(s, 0.7, 1.35, rows, [3.2, 1.35, 7.35], row_h=0.52)
add_card(s, 0.9, 6.0, 11.5, 0.65, "Note", ["For production, WebSocket authentication should be added with a token query parameter or a secure session strategy."], AMBER)
add_footer(s, 13)

prs.save(OUT)
print(OUT)
